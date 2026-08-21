#!/usr/bin/env python3
"""Capability-checked, schema-bound AMOS evaluation toolbox service."""

from __future__ import annotations

import base64
import hashlib
import hmac
import io
import json
import math
import os
import re
import subprocess
import sys
import time
from http.server import BaseHTTPRequestHandler, HTTPServer
from pathlib import Path
from typing import Any, Callable


MAX_REQUEST_BYTES = 50_000_000
TOOLBOX_VERSION = "amos-toolbox-evaluation.v1"
SAFE_NAME = re.compile(r"^[A-Za-z_][A-Za-z0-9_]{0,127}$")

TOOL_AUDIENCES = {
    "spark.dataframe.aggregate.v1": "spark-worker",
    "r.statistics.v1": "r-worker",
    "python.dataframe.aggregate.v1": "python-worker",
    "polars.dataframe.aggregate.v1": "polars-worker",
    "duckdb.readonly.v1": "duckdb-worker",
    "dbt.manifest.validate.v1": "dbt-worker",
    "stats.regression.v1": "stats-worker",
    "stats.forecast.v1": "stats-worker",
    "stats.pca.v1": "stats-worker",
    "spreadsheet.xlsx.v1": "spreadsheet-worker",
    "presentation.pptx.v1": "presentation-worker",
    "notebook.inspect.v1": "notebook-worker",
}

TOOL_OPERATIONS = {
    "spark.dataframe.aggregate.v1": {"aggregate", "query"},
    "r.statistics.v1": {"analyze"},
    "python.dataframe.aggregate.v1": {"aggregate"},
    "polars.dataframe.aggregate.v1": {"aggregate"},
    "duckdb.readonly.v1": {"query"},
    "dbt.manifest.validate.v1": {"validate"},
    "stats.regression.v1": {"analyze"},
    "stats.forecast.v1": {"analyze"},
    "stats.pca.v1": {"analyze"},
    "spreadsheet.xlsx.v1": {"render"},
    "presentation.pptx.v1": {"render"},
    "notebook.inspect.v1": {"inspect"},
}


class ToolboxError(Exception):
    pass


def capability_key() -> bytes:
    path = Path(os.environ.get("AMOS_CAPABILITY_KEY_FILE", "/run/secrets/capability_key"))
    encoded = path.read_text(encoding="ascii").strip()
    if not re.fullmatch(r"[0-9a-f]{64}", encoded):
        raise RuntimeError("capability key must be 32 bytes encoded as lowercase hexadecimal")
    return bytes.fromhex(encoded)


CAPABILITY_KEY = capability_key()


def canonical_json(value: Any) -> bytes:
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"), allow_nan=False).encode("utf-8")


def verify_request(request: dict[str, Any]) -> tuple[str, dict[str, Any], dict[str, Any]]:
    envelope = require_object(request, "capability")
    claims = require_object(envelope, "claims")
    signature = require_string(envelope, "signature")
    expected = hmac.new(CAPABILITY_KEY, canonical_json(claims), hashlib.sha256).digest()
    try:
        supplied = base64.urlsafe_b64decode(signature + "=" * (-len(signature) % 4))
    except Exception as error:
        raise ToolboxError("invalid capability signature encoding") from error
    if not hmac.compare_digest(expected, supplied):
        raise ToolboxError("capability signature verification failed")

    step = require_object(request, "step")
    plan = require_object(request, "plan")
    identity = require_object(request, "identity")
    tool_id = require_string(step, "tool")
    if tool_id not in TOOL_AUDIENCES:
        raise ToolboxError("tool is not implemented by this worker")
    now = int(time.time())
    if claims.get("issuer") != "amos-runtime" or claims.get("audience") != TOOL_AUDIENCES[tool_id]:
        raise ToolboxError("capability issuer or audience mismatch")
    if not isinstance(claims.get("not_before"), int) or not isinstance(claims.get("expires_at"), int):
        raise ToolboxError("capability time window is invalid")
    if now < claims["not_before"] or now >= claims["expires_at"]:
        raise ToolboxError("capability is expired or not active")
    if claims["expires_at"] - claims["not_before"] > 120:
        raise ToolboxError("capability validity window is too large")
    if not isinstance(claims.get("token_id"), str) or not claims["token_id"]:
        raise ToolboxError("capability token identifier is invalid")
    step_hash = "sha256:" + hashlib.sha256(canonical_json(step)).hexdigest()
    bindings = {
        "tenant_id": identity.get("tenant_id"),
        "atxn_id": plan.get("atxn_id"),
        "plan_id": plan.get("plan_id"),
        "step_id": step.get("step_id"),
        "step_hash": step_hash,
        "subject_id": identity.get("subject_id"),
        "tool": tool_id,
        "source_id": step.get("source_id"),
        "policy_epoch": identity.get("policy_epoch"),
        "fencing_token": request.get("fencing_token"),
    }
    for name, value in bindings.items():
        if claims.get(name) != value:
            raise ToolboxError(f"capability {name} binding mismatch")
    if set(claims.get("operations", [])) != TOOL_OPERATIONS[tool_id]:
        raise ToolboxError("capability operations mismatch")
    if plan.get("tenant_id") != identity.get("tenant_id"):
        raise ToolboxError("identity and plan tenant mismatch")
    plan_steps = plan.get("steps")
    if not isinstance(plan_steps, list) or step not in plan_steps:
        raise ToolboxError("invoked step is not present in the approved plan")
    parameters = require_object(step, "parameters")
    relations = parameters.get("relations", [])
    if not isinstance(relations, list) or not all(isinstance(item, str) for item in relations):
        raise ToolboxError("relations must be a string array")
    if set(claims.get("relations", [])) != set(relations):
        raise ToolboxError("capability relation binding mismatch")
    limits = require_object(step, "limits")
    if claims.get("limits") != limits:
        raise ToolboxError("capability limit binding mismatch")
    for name in ("seconds", "rows", "bytes"):
        if not isinstance(limits.get(name), int) or limits[name] <= 0:
            raise ToolboxError("capability limits must be positive integers")
    return tool_id, parameters, limits


def require_object(value: dict[str, Any], name: str) -> dict[str, Any]:
    selected = value.get(name)
    if not isinstance(selected, dict):
        raise ToolboxError(f"{name} must be an object")
    return selected


def require_string(value: dict[str, Any], name: str) -> str:
    selected = value.get(name)
    if not isinstance(selected, str) or not selected:
        raise ToolboxError(f"{name} must be a non-empty string")
    return selected


def rows(parameters: dict[str, Any]) -> list[dict[str, Any]]:
    columns = parameters.get("columns")
    selected = parameters.get("rows")
    if not isinstance(columns, list) or not columns:
        raise ToolboxError("columns must be a non-empty string array")
    safe_columns(columns)
    if (
        not isinstance(selected, list)
        or not selected
        or not all(isinstance(row, list) and len(row) == len(columns) for row in selected)
    ):
        raise ToolboxError("each row must be an array matching the declared columns")
    return [dict(zip(columns, row, strict=True)) for row in selected]


def safe_columns(values: list[str]) -> None:
    if not values or any(not isinstance(value, str) or not SAFE_NAME.fullmatch(value) for value in values):
        raise ToolboxError("column names must be safe identifiers")


def pandas_aggregate(parameters: dict[str, Any]) -> list[dict[str, Any]]:
    import pandas as pd

    frame = pd.DataFrame(rows(parameters))
    group_by = parameters["group_by"]
    safe_columns(group_by)
    named_aggregations: dict[str, Any] = {}
    for metric in parameters["metrics"]:
        function = metric["function"]
        column = metric["column"]
        alias = metric["alias"]
        safe_columns([column, alias])
        if function not in {"count", "sum", "mean", "min", "max"}:
            raise ToolboxError("unsupported dataframe aggregation")
        named_aggregations[alias] = pd.NamedAgg(column=column, aggfunc=function)
    result = frame.groupby(group_by, dropna=False).agg(**named_aggregations).reset_index()
    return json.loads(result.to_json(orient="records"))


def polars_aggregate(parameters: dict[str, Any]) -> list[dict[str, Any]]:
    import polars as pl

    frame = pl.DataFrame(rows(parameters))
    group_by = parameters["group_by"]
    safe_columns(group_by)
    aggregations = []
    for metric in parameters["metrics"]:
        function = metric["function"]
        column = metric["column"]
        alias = metric["alias"]
        safe_columns([column, alias])
        expression = pl.col(column)
        functions_by_name = {
            "count": expression.count,
            "sum": expression.sum,
            "mean": expression.mean,
            "min": expression.min,
            "max": expression.max,
        }
        if function not in functions_by_name:
            raise ToolboxError("unsupported Polars aggregation")
        aggregations.append(functions_by_name[function]().alias(alias))
    return frame.group_by(group_by).agg(aggregations).sort(group_by).to_dicts()


_SPARK_SESSION: Any = None


def spark_aggregate(parameters: dict[str, Any]) -> list[dict[str, Any]]:
    global _SPARK_SESSION
    from pyspark.sql import SparkSession, functions

    if _SPARK_SESSION is None:
        _SPARK_SESSION = (
            SparkSession.builder.master("local[1]")
            .appName("amos-governed-toolbox")
            .config("spark.ui.enabled", "false")
            .config("spark.sql.shuffle.partitions", "1")
            .getOrCreate()
        )
        _SPARK_SESSION.sparkContext.setLogLevel("ERROR")
    frame = _SPARK_SESSION.createDataFrame(rows(parameters))
    for item in parameters["filters"]:
        safe_columns([item["column"]])
        frame = frame.filter(functions.col(item["column"]) == functions.lit(item["equals"]))
    group_by = parameters["group_by"]
    safe_columns(group_by)
    aggregations = []
    for metric in parameters["metrics"]:
        function = metric["function"]
        column = metric["column"]
        alias = metric["alias"]
        safe_columns([column, alias])
        functions_by_name = {
            "count": functions.count,
            "sum": functions.sum,
            "avg": functions.avg,
            "min": functions.min,
            "max": functions.max,
        }
        if function not in functions_by_name:
            raise ToolboxError("unsupported Spark aggregation")
        aggregations.append(functions_by_name[function](column).alias(alias))
    result = frame.groupBy(*group_by).agg(*aggregations).orderBy(*group_by)
    return [row.asDict(recursive=True) for row in result.collect()]


def r_statistics(parameters: dict[str, Any]) -> dict[str, Any]:
    command = ["Rscript", "/opt/amos-toolbox/r_statistics.R"]
    completed = subprocess.run(
        command,
        input=canonical_json(parameters),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
        timeout=60,
    )
    if completed.returncode != 0:
        error = completed.stderr.decode("utf-8", errors="replace")[-2000:]
        raise ToolboxError(f"R worker failed: {error}")
    return json.loads(completed.stdout)


def duckdb_query(parameters: dict[str, Any]) -> list[dict[str, Any]]:
    import duckdb
    import pandas as pd

    query = parameters["sql"].strip()
    normalized = re.sub(r"\s+", " ", query).lower()
    forbidden = re.compile(
        r"\b(attach|call|copy|export|import|install|load|pragma|set|read_[a-z0-9_]*|[a-z0-9_]*_scan)\b"
    )
    if (
        not normalized.startswith("select ")
        or ";" in query
        or "--" in query
        or "/*" in query
        or forbidden.search(normalized)
    ):
        raise ToolboxError("DuckDB accepts exactly one read-only SELECT")
    input_data = pd.DataFrame(rows(parameters))
    connection = duckdb.connect(
        database=":memory:",
        config={
            "enable_external_access": "false",
            "autoload_known_extensions": "false",
            "autoinstall_known_extensions": "false",
        },
    )
    try:
        connection.register("input_data", input_data)
        result = connection.execute(query).fetchdf()
        return json.loads(result.to_json(orient="records"))
    finally:
        connection.close()


def dbt_manifest_validate(parameters: dict[str, Any]) -> dict[str, Any]:
    try:
        manifest = json.loads(parameters["manifest_json"])
    except json.JSONDecodeError as error:
        raise ToolboxError("dbt manifest_json is invalid JSON") from error
    if not isinstance(manifest, dict):
        raise ToolboxError("dbt manifest must be an object")
    nodes = manifest.get("nodes")
    metadata = manifest.get("metadata")
    if not isinstance(nodes, dict) or not isinstance(metadata, dict):
        raise ToolboxError("dbt manifest requires nodes and metadata objects")
    resource_counts: dict[str, int] = {}
    invalid_nodes = []
    for unique_id, node in nodes.items():
        if not isinstance(node, dict) or not isinstance(node.get("resource_type"), str):
            invalid_nodes.append(unique_id)
            continue
        resource = node["resource_type"]
        resource_counts[resource] = resource_counts.get(resource, 0) + 1
    return {
        "valid": not invalid_nodes,
        "node_count": len(nodes),
        "resource_counts": resource_counts,
        "invalid_node_ids": sorted(invalid_nodes),
        "dbt_schema_version": str(metadata.get("dbt_schema_version") or "unknown"),
        "dbt_version": str(metadata.get("dbt_version") or "unknown"),
    }


def regression(parameters: dict[str, Any]) -> dict[str, Any]:
    import numpy as np
    from sklearn.linear_model import LinearRegression
    from sklearn.metrics import mean_squared_error

    x = np.asarray(parameters["x"], dtype=float)
    y = np.asarray(parameters["y"], dtype=float)
    if x.ndim != 2 or y.ndim != 1 or len(x) != len(y) or len(y) < 2:
        raise ToolboxError("regression requires matching two-dimensional x and one-dimensional y")
    model = LinearRegression().fit(x, y)
    predicted = model.predict(x)
    return {
        "intercept": float(model.intercept_),
        "coefficients": [float(value) for value in model.coef_],
        "r_squared": float(model.score(x, y)),
        "rmse": float(math.sqrt(mean_squared_error(y, predicted))),
        "row_count": int(len(y)),
    }


def forecast(parameters: dict[str, Any]) -> dict[str, Any]:
    import numpy as np

    values = np.asarray(parameters["values"], dtype=float)
    horizon = parameters["horizon"]
    if values.ndim != 1 or len(values) < 2:
        raise ToolboxError("forecast requires at least two finite values")
    indices = np.arange(len(values), dtype=float)
    slope, intercept = np.polyfit(indices, values, 1)
    forecast_values = [float(intercept + slope * index) for index in range(len(values), len(values) + horizon)]
    return {
        "method": "linear_trend",
        "slope": float(slope),
        "intercept": float(intercept),
        "forecast": forecast_values,
    }


def pca(parameters: dict[str, Any]) -> dict[str, Any]:
    import numpy as np
    from sklearn.decomposition import PCA

    matrix = np.asarray(parameters["matrix"], dtype=float)
    components = parameters["components"]
    if matrix.ndim != 2 or matrix.shape[0] < 2 or components > min(matrix.shape):
        raise ToolboxError("PCA matrix or component count is invalid")
    model = PCA(n_components=components, svd_solver="full").fit(matrix)
    transformed = model.transform(matrix)
    return {
        "components": model.components_.tolist(),
        "explained_variance": model.explained_variance_.tolist(),
        "explained_variance_ratio": model.explained_variance_ratio_.tolist(),
        "transformed": transformed.tolist(),
    }


def spreadsheet(parameters: dict[str, Any]) -> dict[str, Any]:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    for sheet_spec in parameters["sheets"]:
        title = sheet_spec["title"]
        if len(title) > 31 or any(character in title for character in "[]:*?/\\"):
            raise ToolboxError("worksheet title is invalid")
        sheet = workbook.create_sheet(title)
        for row in sheet_spec["rows"]:
            sheet.append(row)
        sheet.freeze_panes = "A2" if sheet_spec.get("freeze_header") else None
    buffer = io.BytesIO()
    workbook.save(buffer)
    body = buffer.getvalue()
    return encoded_artifact(body, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx")


def presentation(parameters: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    for slide_spec in parameters["slides"]:
        layout = deck.slide_layouts[1]
        slide = deck.slides.add_slide(layout)
        slide.shapes.title.text = slide_spec["title"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(slide_spec["bullets"]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
    buffer = io.BytesIO()
    deck.save(buffer)
    body = buffer.getvalue()
    return encoded_artifact(body, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx")


def notebook_inspect(parameters: dict[str, Any]) -> dict[str, Any]:
    import nbformat

    try:
        notebook = nbformat.reads(parameters["notebook_json"], as_version=4)
    except Exception as error:
        raise ToolboxError("notebook_json is not a valid version 4 notebook") from error
    cell_counts: dict[str, int] = {}
    code_hashes = []
    for cell in notebook.cells:
        cell_counts[cell.cell_type] = cell_counts.get(cell.cell_type, 0) + 1
        if cell.cell_type == "code":
            code_hashes.append("sha256:" + hashlib.sha256(cell.source.encode("utf-8")).hexdigest())
    return {
        "format": f"{notebook.nbformat}.{notebook.nbformat_minor}",
        "cell_count": len(notebook.cells),
        "cell_counts": cell_counts,
        "code_hashes": code_hashes,
        "executed": False,
    }


def encoded_artifact(body: bytes, media_type: str, extension: str) -> dict[str, Any]:
    return {
        "media_type": media_type,
        "extension": extension,
        "bytes": len(body),
        "sha256": hashlib.sha256(body).hexdigest(),
        "base64": base64.b64encode(body).decode("ascii"),
    }


EXECUTORS: dict[str, Callable[[dict[str, Any]], Any]] = {
    "spark.dataframe.aggregate.v1": spark_aggregate,
    "r.statistics.v1": r_statistics,
    "python.dataframe.aggregate.v1": pandas_aggregate,
    "polars.dataframe.aggregate.v1": polars_aggregate,
    "duckdb.readonly.v1": duckdb_query,
    "dbt.manifest.validate.v1": dbt_manifest_validate,
    "stats.regression.v1": regression,
    "stats.forecast.v1": forecast,
    "stats.pca.v1": pca,
    "spreadsheet.xlsx.v1": spreadsheet,
    "presentation.pptx.v1": presentation,
    "notebook.inspect.v1": notebook_inspect,
}


def execute(request: dict[str, Any]) -> dict[str, Any]:
    tool_id, parameters, limits = verify_request(request)
    started = time.monotonic()
    output = EXECUTORS[tool_id](parameters)
    elapsed = time.monotonic() - started
    if elapsed > limits["seconds"]:
        raise ToolboxError("tool execution exceeded its time limit")
    encoded = canonical_json(output)
    if len(encoded) > limits["bytes"]:
        raise ToolboxError("tool output exceeded its byte limit")
    row_count = len(output) if isinstance(output, list) else 1
    if row_count > limits["rows"]:
        raise ToolboxError("tool output exceeded its row limit")
    return {
        "status": "pass",
        "tool_id": tool_id,
        "toolbox_version": TOOLBOX_VERSION,
        "runtime_versions": runtime_versions(tool_id),
        "output": output,
        "row_count": row_count,
        "byte_count": len(encoded),
        "latency_ms": int(elapsed * 1000),
    }


def runtime_versions(tool_id: str) -> dict[str, str]:
    versions = {"python": sys.version.split()[0]}
    if tool_id.startswith("spark."):
        import pyspark
        versions["pyspark"] = pyspark.__version__
    elif tool_id.startswith("r."):
        completed = subprocess.run(["Rscript", "--version"], capture_output=True, check=False, timeout=5)
        version = (completed.stderr or completed.stdout).decode("utf-8", errors="replace").strip()
        versions["r"] = version or "unknown"
    elif tool_id.startswith("python."):
        import pandas
        versions["pandas"] = pandas.__version__
    elif tool_id.startswith("polars."):
        import polars
        versions["polars"] = polars.__version__
    elif tool_id.startswith("duckdb."):
        import duckdb
        versions["duckdb"] = duckdb.__version__
    elif tool_id.startswith("stats."):
        import numpy
        import sklearn
        versions.update({"numpy": numpy.__version__, "scikit_learn": sklearn.__version__})
    elif tool_id.startswith("spreadsheet."):
        import openpyxl
        versions["openpyxl"] = openpyxl.__version__
    elif tool_id.startswith("presentation."):
        import pptx
        versions["python_pptx"] = pptx.__version__
    elif tool_id.startswith("notebook."):
        import nbformat
        versions["nbformat"] = nbformat.__version__
    return versions


class Handler(BaseHTTPRequestHandler):
    server_version = "amos-toolbox/1"

    def do_GET(self) -> None:
        if self.path != "/health":
            self.respond(404, {"error": "not found"})
            return
        self.respond(200, {"status": "ok", "version": TOOLBOX_VERSION, "tools": sorted(EXECUTORS)})

    def do_POST(self) -> None:
        if self.path != "/execute":
            self.respond(404, {"error": "not found"})
            return
        try:
            length = int(self.headers.get("content-length", "0"))
            if length <= 0 or length > MAX_REQUEST_BYTES:
                raise ToolboxError("request size is invalid")
            request = json.loads(self.rfile.read(length))
            if not isinstance(request, dict):
                raise ToolboxError("request must be an object")
            self.respond(200, execute(request))
        except ToolboxError as error:
            self.respond(422, {"error": str(error)})
        except Exception:
            self.respond(500, {"error": "toolbox execution failed"})

    def respond(self, status: int, value: dict[str, Any]) -> None:
        body = canonical_json(value)
        self.send_response(status)
        self.send_header("content-type", "application/json")
        self.send_header("content-length", str(len(body)))
        self.send_header("cache-control", "no-store")
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, format: str, *args: Any) -> None:
        sys.stderr.write("amos-toolbox: " + format % args + "\n")


if __name__ == "__main__":
    address = os.environ.get("AMOS_TOOLBOX_BIND", "0.0.0.0")
    port = int(os.environ.get("AMOS_TOOLBOX_PORT", "9000"))
    HTTPServer((address, port), Handler).serve_forever()
